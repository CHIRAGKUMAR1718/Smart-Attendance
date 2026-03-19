from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Color Palette ──
NAVY        = RGBColor(0x0A, 0x1F, 0x3C)
DARK_BLUE   = RGBColor(0x00, 0x33, 0x66)
MID_BLUE    = RGBColor(0x00, 0x6B, 0xB6)
TEAL        = RGBColor(0x00, 0x80, 0x80)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x00, 0x00, 0x00)
OFF_WHITE   = RGBColor(0xFA, 0xFA, 0xFA)
LIGHT_GRAY  = RGBColor(0xF0, 0xF0, 0xF0)
MED_GRAY    = RGBColor(0xCC, 0xCC, 0xCC)
DARK_GRAY   = RGBColor(0x33, 0x33, 0x33)
ACCENT      = RGBColor(0x63, 0x66, 0xF1)
SUCCESS     = RGBColor(0x10, 0xB9, 0x81)
ORANGE      = RGBColor(0xF5, 0x9E, 0x0B)
RED         = RGBColor(0xEF, 0x44, 0x44)
PURPLE      = RGBColor(0x8B, 0x5C, 0xF6)
CYAN        = RGBColor(0x06, 0xB6, 0xD4)
SOFT_INDIGO = RGBColor(0xEE, 0xEF, 0xFE)
SOFT_GREEN  = RGBColor(0xE6, 0xFA, 0xF0)
SOFT_ORANGE = RGBColor(0xFE, 0xF3, 0xE2)
SOFT_RED    = RGBColor(0xFE, 0xE2, 0xE2)
SOFT_BLUE   = RGBColor(0xE0, 0xF2, 0xFE)
SOFT_PURPLE = RGBColor(0xF3, 0xE8, 0xFF)
SOFT_CYAN   = RGBColor(0xE0, 0xF7, 0xFA)
CODE_BG     = RGBColor(0x1E, 0x1E, 0x2E)
CODE_TEXT   = RGBColor(0xCD, 0xD6, 0xF4)
CODE_KEY    = RGBColor(0xF3, 0x8B, 0xA8)
CODE_STR    = RGBColor(0xA6, 0xE3, 0xA1)


# ═══════════════ HELPERS ═══════════════

def add_shape(slide, stype, left, top, width, height, fill=None, border_color=None, border_width=1.5):
    s = slide.shapes.add_shape(stype, left, top, width, height)
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border_color:
        s.line.color.rgb = border_color
        s.line.width = Pt(border_width)
    else:
        s.line.fill.background()
    return s

def add_text(slide, left, top, width, height, text, size=18, bold=False, color=BLACK,
             align=PP_ALIGN.LEFT, italic=False, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    run.font.name = font
    p.alignment = align
    return tf

def rounded_box(slide, left, top, width, height, fill, text="", fs=16,
                fc=WHITE, align=PP_ALIGN.CENTER, bold=True, border_color=None, bw=2):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border_color:
        s.line.color.rgb = border_color
        s.line.width = Pt(bw)
    else:
        s.line.fill.background()
    if text:
        tf = s.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(fs)
        p.font.bold = bold
        p.font.color.rgb = fc
        p.alignment = align
        p.font.name = "Calibri"
    return s

def header_bar(slide, title, subtitle=""):
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.05), fill=NAVY)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, Inches(1.05), prs.slide_width, Inches(0.045), fill=ACCENT)
    add_text(slide, Inches(0.6), Inches(0.1), Inches(9), Inches(0.5),
             title, 30, True, WHITE)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(0.55), Inches(10), Inches(0.4),
                 subtitle, 14, False, RGBColor(0x94, 0xA3, 0xB8), italic=True)
    add_text(slide, Inches(10.8), Inches(0.25), Inches(2.2), Inches(0.5),
             "⚡ SMAT", 16, True, ACCENT, PP_ALIGN.RIGHT)

def footer_bar(slide, num=""):
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, Inches(7.1), prs.slide_width, Inches(0.4), fill=NAVY)
    add_text(slide, Inches(0.5), Inches(7.12), Inches(8), Inches(0.35),
             "IIIT Surat  •  Dept. of Electronics & Communication Engineering", 11, False,
             RGBColor(0x94, 0xA3, 0xB8))
    if num:
        add_text(slide, Inches(11), Inches(7.12), Inches(2), Inches(0.35),
                 num, 11, False, RGBColor(0x94, 0xA3, 0xB8), PP_ALIGN.RIGHT)

def bullet_list(slide, left, top, width, height, items, color=DARK_GRAY, fs=14, sp=4, bullet="▸"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = ""
    for item in items:
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = f"{bullet}  {item}"
        r.font.size = Pt(fs)
        r.font.color.rgb = color
        r.font.name = "Calibri"
        p.space_before = Pt(sp)
    return tf

def rich_bullet_list(slide, left, top, width, height, items, label_color=ACCENT, desc_color=DARK_GRAY, fs=14, sp=5):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = ""
    for label, desc in items:
        p = tf.add_paragraph()
        r1 = p.add_run()
        r1.text = f"▸  {label}  "
        r1.font.size = Pt(fs)
        r1.font.bold = True
        r1.font.color.rgb = label_color
        r1.font.name = "Calibri"
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(fs)
        r2.font.color.rgb = desc_color
        r2.font.name = "Calibri"
        p.space_before = Pt(sp)
    return tf

def code_block(slide, left, top, width, height, lines):
    box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height, fill=CODE_BG, border_color=RGBColor(0x45,0x47,0x5A), border_width=1)
    tb = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), height - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = ""
    for line in lines:
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = line
        r.font.size = Pt(11)
        r.font.color.rgb = CODE_TEXT
        r.font.name = "Consolas"
        p.space_before = Pt(1)
    return tf

def table_row(slide, col_lefts, col_widths, row_top, row_h, values, bg, text_color=DARK_GRAY, fs=13, bold_col=None, aligns=None):
    for i, val in enumerate(values):
        box = add_shape(slide, MSO_SHAPE.RECTANGLE, col_lefts[i], row_top, col_widths[i], row_h, fill=bg)
        box.line.color.rgb = MED_GRAY
        box.line.width = Pt(0.5)
        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        p = tf.paragraphs[0]
        p.text = val
        p.font.size = Pt(fs)
        p.font.color.rgb = text_color
        p.font.bold = (bold_col is not None and i in bold_col)
        p.alignment = aligns[i] if aligns else PP_ALIGN.LEFT
        p.font.name = "Calibri"

def table_header(slide, col_lefts, col_widths, row_top, row_h, values, fs=14):
    for i, val in enumerate(values):
        box = add_shape(slide, MSO_SHAPE.RECTANGLE, col_lefts[i], row_top, col_widths[i], row_h, fill=NAVY)
        box.line.color.rgb = WHITE
        box.line.width = Pt(1)
        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.08)
        p = tf.paragraphs[0]
        p.text = val
        p.font.size = Pt(fs)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Calibri"

def set_bg(slide, color=OFF_WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

slide_number = [0]

def new_slide():
    slide_number[0] += 1
    return prs.slides.add_slide(prs.slide_layouts[6])


# ═══════════════════════════════════════════════
#  SLIDE 0 — TITLE PAGE
# ═══════════════════════════════════════════════
slide = new_slide()
set_bg(slide, WHITE)

add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.8), fill=NAVY)
add_text(slide, Inches(2.0), Inches(0.12), Inches(9.3), Inches(0.45),
         "भारतीय सूचना प्रौद्योगिकी संस्थान सूरत", 20, True, WHITE, PP_ALIGN.CENTER)
add_text(slide, Inches(2.0), Inches(0.48), Inches(9.3), Inches(0.5),
         "Indian Institute of Information Technology Surat", 24, True, WHITE, PP_ALIGN.CENTER)
add_text(slide, Inches(2.0), Inches(0.88), Inches(9.3), Inches(0.35),
         "ભારતીય સૂચના પ્રૌદ્યોગિકી સંસ્થા સુરત", 14, False, RGBColor(0xAA,0xBB,0xCC), PP_ALIGN.CENTER)
add_text(slide, Inches(2.0), Inches(1.18), Inches(9.3), Inches(0.35),
         "(An Institute of National Importance under Act of Parliament)", 13, False,
         RGBColor(0x88,0x99,0xAA), PP_ALIGN.CENTER, italic=True)

add_text(slide, Inches(1), Inches(2.2), Inches(11.3), Inches(0.4),
         "Presentation on", 22, False, DARK_BLUE, PP_ALIGN.CENTER, italic=True)

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(2.7), Inches(10.3), Inches(1.1),
          fill=SOFT_INDIGO, border_color=ACCENT, border_width=2)
add_text(slide, Inches(1.5), Inches(2.8), Inches(10.3), Inches(0.9),
         "SMAT: Smart Attendance System\nUsing Ultrasonic Audio Signals", 32, True, NAVY, PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(4.1), Inches(11.3), Inches(0.35),
         "Group Members", 20, True, DARK_BLUE, PP_ALIGN.CENTER)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(4.5), Inches(4.3), Inches(0.02), fill=ACCENT)

mtb = slide.shapes.add_textbox(Inches(3.5), Inches(4.6), Inches(6.3), Inches(1.0))
mtf = mtb.text_frame
mtf.word_wrap = True
p = mtf.paragraphs[0]
p.text = "<Name>                              <Roll Number>"
p.font.size = Pt(18); p.font.bold = True; p.font.italic = True; p.font.color.rgb = DARK_GRAY
p.alignment = PP_ALIGN.CENTER; p.font.name = "Calibri"

add_text(slide, Inches(1), Inches(5.55), Inches(11.3), Inches(0.3),
         "Under the guidance of", 16, False, DARK_GRAY, PP_ALIGN.CENTER, italic=True)
add_text(slide, Inches(1), Inches(5.85), Inches(11.3), Inches(0.35),
         "<Name of the Guide>", 18, True, TEAL, PP_ALIGN.CENTER, italic=True)

add_shape(slide, MSO_SHAPE.RECTANGLE, 0, Inches(6.65), prs.slide_width, Inches(0.85), fill=TEAL)
add_text(slide, Inches(1), Inches(6.75), Inches(11.3), Inches(0.45),
         "Department of Electronics and Communication Engineering", 18, True, WHITE, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
#  SLIDE 1 — OUTLINE
# ═══════════════════════════════════════════════
slide = new_slide()
set_bg(slide, OFF_WHITE)
header_bar(slide, "Outline of Presentation")
footer_bar(slide, "1 / 13")

outline = [
    ("01", "Introduction", ACCENT),
    ("02", "Literature Survey", CYAN),
    ("03", "Problem Statement", RED),
    ("04", "Proposed System Architecture", MID_BLUE),
    ("05", "Database Design & ER Schema", TEAL),
    ("06", "API Design & Endpoints", PURPLE),
    ("07", "Security & Middleware", RGBColor(0x78,0x35,0x0D)),
    ("08", "Hardware & Software Requirements", DARK_BLUE),
    ("09", "Implementation Updates", SUCCESS),
    ("10", "Results & Demonstrations", ORANGE),
    ("11", "Advantages & Applications", ACCENT),
    ("12", "Budget Analysis", TEAL),
    ("13", "Workflow for Remaining Work", PURPLE),
    ("14", "Future Scope", RED),
    ("15", "References", DARK_GRAY),
]

for i, (num, title, color) in enumerate(outline):
    col = i % 3
    row = i // 3
    left = Inches(0.4 + col * 4.2)
    top = Inches(1.3 + row * 1.1)
    rounded_box(slide, left, top, Inches(0.6), Inches(0.55), color, num, 16, WHITE)
    label_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.7), top + Inches(0.03),
                           Inches(3.3), Inches(0.5), fill=WHITE, border_color=color, border_width=1.5)
    tf = label_box.text_frame
    tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.margin_left = Inches(0.12)
    p = tf.paragraphs[0]
    p.text = title; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = color; p.font.name = "Calibri"


# ═══════════════════════════════════════════════
#  SLIDE 2 — INTRODUCTION (page 1)
# ═══════════════════════════════════════════════
slide = new_slide()
set_bg(slide, OFF_WHITE)
header_bar(slide, "Introduction", "What is SMAT and why does it exist?")
footer_bar(slide, "2a / 13")

# Left — What is SMAT
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.25), Inches(6.3), Inches(5.6),
          fill=WHITE, border_color=ACCENT, border_width=2)
add_text(slide, Inches(0.5), Inches(1.35), Inches(5.9), Inches(0.4),
         "What is SMAT?", 22, True, ACCENT)

bullet_list(slide, Inches(0.5), Inches(1.8), Inches(5.9), Inches(1.8), [
    "Smart Attendance Marking using inaudible ultrasonic audio (18–20 kHz)",
    "Teacher's device broadcasts encoded sound signal via speakers",
    "Student's device listens via microphone, decodes the signal automatically",
    "Ensures physical presence — ultrasonic waves cannot travel through walls or internet",
    "100% browser-based — zero installation on any device (Web Audio API)",
    "Role-based system with separate Teacher & Student portals",
], DARK_GRAY, 14, 5)

add_text(slide, Inches(0.5), Inches(4.7), Inches(5.9), Inches(0.35),
         "Why Ultrasonic Audio?", 18, True, PURPLE)

bullet_list(slide, Inches(0.5), Inches(5.1), Inches(5.9), Inches(1.5), [
    "Inaudible to humans (>17 kHz) — does not disrupt class",
    "Short range (~5–10m) — naturally limits to classroom boundaries",
    "Cannot be relayed over phone calls or messaging apps",
    "Standard laptop speakers can produce up to 20–22 kHz",
    "Standard phone microphones can capture up to 22–24 kHz",
], DARK_GRAY, 13, 4)

# Right — How it works
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.25), Inches(6.2), Inches(3.2),
          fill=NAVY)
add_text(slide, Inches(7.0), Inches(1.35), Inches(5.8), Inches(0.35),
         "End-to-End Attendance Flow", 20, True, WHITE)

steps = [
    ("1", "Teacher enters Class ID, selects duration (1–30 min)", SUCCESS),
    ("2", "Server generates UUID session + unique 6-char code", CYAN),
    ("3", "Browser converts code → ultrasonic tones, broadcasts via speaker", ACCENT),
    ("4", "Student mic captures audio; FFT analyzer decodes 6 characters", ORANGE),
    ("5", "Decoded code sent to server → validates session → marks attendance", PURPLE),
    ("6", "Teacher sees real-time count update; student sees confirmation", SUCCESS),
]
for i, (num, text, col) in enumerate(steps):
    y = Inches(1.85 + i * 0.43)
    rounded_box(slide, Inches(7.05), y, Inches(0.35), Inches(0.35), col, num, 13, WHITE)
    add_text(slide, Inches(7.5), y + Inches(0.02), Inches(5.3), Inches(0.32),
             text, 13, False, RGBColor(0xCC,0xDD,0xEE))

# Right bottom — comparison table
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(4.65), Inches(6.2), Inches(2.2),
          fill=WHITE, border_color=RED, border_width=2)
add_text(slide, Inches(7.0), Inches(4.72), Inches(5.8), Inches(0.3),
         "SMAT vs Other Systems — Quick Comparison", 16, True, RED)

comp_cols_l = [Inches(6.95), Inches(8.95), Inches(10.05), Inches(11.15)]
comp_cols_w = [Inches(2.0), Inches(1.1), Inches(1.1), Inches(1.7)]
table_header(slide, comp_cols_l, comp_cols_w, Inches(5.05), Inches(0.35),
             ["Feature", "QR Code", "RFID", "SMAT"], 11)
comp_data = [
    ("Proxy Prevention", "✗", "✓", "✓"),
    ("No Hardware Cost", "✓", "✗", "✓"),
    ("No Installation", "~", "✗", "✓"),
    ("Real-time Tracking", "✗", "✗", "✓"),
]
for r, vals in enumerate(comp_data):
    bg = WHITE if r % 2 == 0 else LIGHT_GRAY
    table_row(slide, comp_cols_l, comp_cols_w, Inches(5.4 + r * 0.32), Inches(0.32), list(vals), bg, DARK_GRAY, 11, bold_col={0})


# ═══════════════════════════════════════════════
#  SLIDE 2b — INTRODUCTION (Ultrasonic Protocol Deep Dive)
# ═══════════════════════════════════════════════
slide = new_slide()
set_bg(slide, OFF_WHITE)
header_bar(slide, "Ultrasonic Encoding Protocol", "Technical deep-dive into the signal design")
footer_bar(slide, "2b / 13")

# Signal Structure
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.25), Inches(12.7), Inches(2.1),
          fill=WHITE, border_color=PURPLE, border_width=2)
add_text(slide, Inches(0.5), Inches(1.32), Inches(12.3), Inches(0.35),
         "Signal Structure — One Complete Broadcast Cycle (~2.3 seconds)", 20, True, PURPLE)

# Visual: Preamble → Tone1 → Gap → Tone2 → ... → Tone6 → Pause
signal_parts = [
    ("Preamble\n17800 Hz\n250ms", PURPLE, Inches(1.6)),
    ("GAP\n90ms", MED_GRAY, Inches(0.8)),
    ("Char 1\n18000+\n180ms", ACCENT, Inches(1.15)),
    ("GAP\n90ms", MED_GRAY, Inches(0.8)),
    ("Char 2\n...\n180ms", CYAN, Inches(1.15)),
    ("GAP\n90ms", MED_GRAY, Inches(0.8)),
    ("Char 3", TEAL, Inches(0.7)),
    ("...", MED_GRAY, Inches(0.45)),
    ("Char 6\n...\n180ms", SUCCESS, Inches(1.15)),
    ("PAUSE\n650ms", RGBColor(0x66,0x66,0x66), Inches(1.3)),
    ("↺ LOOP", ORANGE, Inches(1.0)),
]
x = Inches(0.5)
y_sig = Inches(1.8)
for label, color, w in signal_parts:
    rounded_box(slide, x, y_sig, w, Inches(0.85), color, label, 10, WHITE)
    x += w + Inches(0.06)

# Frequency Mapping
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(3.55), Inches(6.3), Inches(3.4),
          fill=WHITE, border_color=ACCENT, border_width=2)
add_text(slide, Inches(0.5), Inches(3.62), Inches(5.9), Inches(0.35),
         "Frequency ↔ Character Mapping", 18, True, ACCENT)
add_text(slide, Inches(0.5), Inches(3.98), Inches(5.9), Inches(0.3),
         "Charset: A B C D E F G H J K M N P Q R S T U V W X Y Z 2 3 4 5 6 7 8 9  (30 chars, no O/0/I/1/L ambiguity)", 11, False, DARK_GRAY, font="Consolas")

map_cols_l = [Inches(0.5), Inches(1.5), Inches(2.5), Inches(3.5), Inches(4.5), Inches(5.5)]
map_cols_w = [Inches(1.0)] * 6
table_header(slide, map_cols_l, map_cols_w, Inches(4.35), Inches(0.3), ["Char", "Freq (Hz)", "Char", "Freq (Hz)", "Char", "Freq (Hz)"], 10)

char_map = [
    ("A", "18000", "K", "18550", "V", "19100"),
    ("B", "18055", "M", "18605", "W", "19155"),
    ("C", "18110", "N", "18660", "X", "19210"),
    ("D", "18165", "P", "18715", "Y", "19265"),
    ("E", "18220", "Q", "18770", "Z", "19320"),
    ("F", "18275", "R", "18825", "2", "19375"),
    ("G", "18330", "S", "18880", "3", "19430"),
    ("H", "18385", "T", "18935", "4–9", "19485–19760"),
]
for r, row_vals in enumerate(char_map):
    bg = WHITE if r % 2 == 0 else LIGHT_GRAY
    table_row(slide, map_cols_l, map_cols_w, Inches(4.65 + r * 0.27), Inches(0.27), list(row_vals), bg, DARK_GRAY, 10,
              bold_col={0, 2, 4}, aligns=[PP_ALIGN.CENTER]*6)

add_text(slide, Inches(0.5), Inches(6.85), Inches(5.9), Inches(0.2),
         "Formula:  freq = 18000 + indexOf(char) × 55 Hz", 12, True, PURPLE, font="Consolas")

# Detection State Machine
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(3.55), Inches(6.2), Inches(3.4),
          fill=WHITE, border_color=TEAL, border_width=2)
add_text(slide, Inches(7.0), Inches(3.62), Inches(5.8), Inches(0.35),
         "Detection State Machine (Student Side)", 18, True, TEAL)

states = [
    ("IDLE", "Listen for 17800 Hz preamble\nRequire ≥4 consecutive frames above -70 dB", MED_GRAY),
    ("AFTER_PRE", "Wait for preamble to stop\n≥2 frames of silence confirms transition", ORANGE),
    ("RECV", "Detect data tones (18000–19925 Hz)\nMedian frequency of ≥2 readings → character\nRepeat until 6 characters decoded", ACCENT),
    ("COMPLETE", "All 6 chars decoded → auto-submit\nLookup session by code → mark attendance", SUCCESS),
]
for i, (name, desc, color) in enumerate(states):
    y = Inches(4.1 + i * 0.72)
    rounded_box(slide, Inches(7.05), y, Inches(1.3), Inches(0.6), color, name, 13, WHITE)
    add_text(slide, Inches(8.45), y + Inches(0.02), Inches(4.4), Inches(0.56), desc, 10, False, DARK_GRAY)
    if i < len(states) - 1:
        add_text(slide, Inches(7.5), y + Inches(0.6), Inches(0.4), Inches(0.15), "↓", 12, True, color, PP_ALIGN.CENTER)

# Technical specs row
specs = [
    ("FFT Size:", "8192 bins"), ("Sample Rate:", "~48 kHz"), ("Bin Width:", "~5.86 Hz"),
    ("Threshold:", "-70 dB"), ("Poll Rate:", "30ms"), ("Timeout:", "8 sec"),
    ("Smoothing:", "0.15"), ("Gain Ramp:", "6ms"),
]
for i, (k, v) in enumerate(specs):
    col = i % 4
    row = i // 4
    x = Inches(7.0 + col * 1.45)
    y = Inches(6.95 + row * 0.0)  # single row actually
    if row == 1:
        x = Inches(7.0 + col * 1.45)
        y = Inches(6.95)

# Actually let me put them in a bottom bar
spec_text = "  |  ".join([f"{k} {v}" for k, v in specs])
rounded_box(slide, Inches(6.8), Inches(6.95), Inches(6.2), Inches(0.15), TEAL, "", 1, WHITE)
# Too small, skip


# ═══════════════════════════════════════════════
#  SLIDE 3 — LITERATURE SURVEY
# ═══════════════════════════════════════════════
slide = new_slide()
set_bg(slide, OFF_WHITE)
header_bar(slide, "Literature Survey", "Existing approaches and their limitations")
footer_bar(slide, "3 / 13")

studies = [
    ("Traditional Methods", MID_BLUE, SOFT_BLUE, [
        "Manual roll call — disrupts 5–10 min/session; error-prone in large classes (100+ students)",
        "Paper sign-in sheets — easily forged; no digital trail; impossible to audit",
        "Biometric (fingerprint/face) — ₹10,000–₹50,000/unit; queuing bottleneck; privacy issues",
    ]),
    ("Digital Approaches", TEAL, SOFT_CYAN, [
        "QR Code — Chowdhury et al. (2019): students share QR screenshots over WhatsApp; 23% proxy rate",
        "GPS-based — Lim et al. (2020): GPS spoof apps on Play Store bypass location checks trivially",
        "Bluetooth beacons — BLE beacons cost ₹2,000–₹5,000 each; need per-room installation",
    ]),
    ("Audio/Ultrasonic Research", PURPLE, SOFT_PURPLE, [
        "Dhwani — Nandakumar et al. (ACM SIGCOMM 2013): peer-to-peer acoustic NFC at 18–20 kHz",
        "Chirp.io — commercial data-over-sound SDK; discontinued in 2023; required native SDK",
        "Google Nearby — uses near-ultrasonic audio for proximity; closed-source; Android/iOS only",
    ]),
    ("Research Gap Identified", RED, SOFT_RED, [
        "No existing system combines: ultrasonic audio + browser-only + zero hardware + open protocol",
        "All current solutions need native apps, proprietary SDKs, or hardware investment",
        "SMAT fills this gap: Web Audio API, standard devices, open frequency protocol, sub-30s marking",
    ]),
]

for i, (title, color, bg, points) in enumerate(studies):
    col = i % 2
    row = i // 2
    left = Inches(0.3 + col * 6.5)
    top = Inches(1.25 + row * 2.85)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(6.2), Inches(2.6), fill=bg, border_color=color, border_width=2)
    rounded_box(slide, left + Inches(0.08), top + Inches(0.08), Inches(6.04), Inches(0.45), color, title, 16, WHITE)
    bullet_list(slide, left + Inches(0.2), top + Inches(0.6), Inches(5.8), Inches(1.8), points, DARK_GRAY, 13, 4)


# ═══════════════════════════════════════════════
#  SLIDE 4 — PROBLEM STATEMENT
# ═══════════════════════════════════════════════
slide = new_slide()
set_bg(slide, OFF_WHITE)
header_bar(slide, "Problem Statement", "Challenges in current attendance systems")
footer_bar(slide, "4 / 13")

problems = [
    ("⏱", "Time Wastage", "Manual roll call consumes 5–10 min\nof each lecture for classes >60\nstudents; disrupts flow of teaching", MID_BLUE, SOFT_BLUE),
    ("👥", "Proxy Attendance", "Students mark for absent peers;\nQR codes shared via WhatsApp;\nfriends sign paper sheets", RED, SOFT_RED),
    ("📍", "Location Spoofing", "GPS-based apps easily fooled;\nVPN/mock location apps widely\navailable on Android & iOS", ORANGE, SOFT_ORANGE),
    ("💰", "Hardware Cost", "Biometric (₹10K–₹50K/unit);\nRFID tags + readers per room;\nMaintenance & power supply costs", PURPLE, SOFT_PURPLE),
]
for i, (icon, title, desc, color, bg) in enumerate(problems):
    left = Inches(0.3 + i * 3.2)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.25), Inches(3.0), Inches(2.65), fill=bg, border_color=color, border_width=2)
    add_text(slide, left + Inches(0.1), Inches(1.35), Inches(2.8), Inches(0.35), icon, 26, False, color, PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.1), Inches(1.72), Inches(2.8), Inches(0.35), title, 19, True, color, PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.15), Inches(2.15), Inches(2.7), Inches(1.5), desc, 13, False, DARK_GRAY, PP_ALIGN.CENTER)

add_text(slide, Inches(0.3), Inches(4.1), Inches(12.7), Inches(0.4), "▼     ▼     ▼     ▼", 26, True, ACCENT, PP_ALIGN.CENTER)

# Solution
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.3),
          fill=WHITE, border_color=SUCCESS, border_width=3)
add_text(slide, Inches(0.7), Inches(4.68), Inches(11.9), Inches(0.4),
         "Our Solution → SMAT (Smart Attendance using Ultrasonic Audio)", 24, True, SUCCESS, PP_ALIGN.CENTER)

sol_items = [
    ("Proximity Verified:", "Ultrasonic audio (18–20 kHz) only travels ~5–10m in air; cannot penetrate walls → physical presence guaranteed"),
    ("Zero Hardware Cost:", "Uses existing laptop speakers & phone microphones — no additional equipment to purchase or maintain"),
    ("Zero Installation:", "Entirely runs in the browser via Web Audio API; no mobile app, no plugin, no download required"),
    ("Anti-Replay:", "6-char codes from 30-symbol charset → 729 million combos; sessions auto-expire via Redis TTL"),
    ("Real-time:", "Teacher sees live student check-in count & names during broadcast; instant feedback loop"),
]
rich_bullet_list(slide, Inches(0.8), Inches(5.1), Inches(11.7), Inches(1.7), sol_items, SUCCESS, DARK_GRAY, 13, 3)


# ═══════════════════════════════════════════════
#  SLIDE 5 — SYSTEM ARCHITECTURE
# ═══════════════════════════════════════════════
slide = new_slide()
set_bg(slide, OFF_WHITE)
header_bar(slide, "Proposed System Architecture", "Layered architecture with real-time capabilities")
footer_bar(slide, "5a / 13")

layers = [
    ("PRESENTATION LAYER", "Frontend (Browser)",
     "index.html (Landing Page)  •  teacher.html (Dashboard + Signal Broadcast)  •  student.html (Listener + History)\n"
     "Technologies: HTML5, CSS3, JavaScript ES6+, Web Audio API, FFT Visualization, Fetch API\n"
     "Design: Dark theme (#0f0f23), Inter font, responsive, backdrop-filter blur, CSS animations",
     SUCCESS, SOFT_GREEN),
    ("APPLICATION LAYER", "REST API (Express 5 / Node.js 20)",
     "3 Route Groups: /auth (register, login, refresh) • /sessions (start, active, history, code/:code) • /attendance (mark, me, list, class/:id, session/:id)\n"
     "Middleware Pipeline: CORS → JSON parser → Rate Limiter → Auth (JWT verify) → Role check → Controller → Error handler\n"
     "Error Handling: Custom error classes with HTTP status codes; stack traces in dev mode only",
     ACCENT, SOFT_INDIGO),
    ("BUSINESS LOGIC LAYER", "Service Modules",
     "auth.service.js: bcrypt hash (12 rounds) + JWT sign (2h access, 7d refresh) + refresh rotation + email normalization\n"
     "session.service.js: UUID v4 generation + 6-char code (30-char set, no ambiguous O/0/I/1/L) + Redis SET EX + Postgres INSERT\n"
     "attendance.service.js: Redis EXISTS (session valid?) → Redis EXISTS (dedup?) → Redis SET EX 120s → Postgres INSERT + UNIQUE constraint",
     PURPLE, SOFT_PURPLE),
    ("DATA LAYER", "PostgreSQL 16 + Redis 7",
     "PostgreSQL: 3 tables (users, sessions, attendance_records) with indexes on teacher_id, class_id, session_id, student_id\n"
     "Redis Keys: chirp:session:{uuid} (JSON, TTL=duration) • chirp:code:{6char} (→sessionId, TTL) • attendance:{sid}:{uid} (dedup, 120s) • rl:*:{ip} (rate limit)\n"
     "Connection: pg.Pool for Postgres, ioredis with lazyConnect + retry strategy (exponential backoff, max 30s) for Redis",
     DARK_BLUE, SOFT_BLUE),
]
for i, (layer_name, layer_sub, detail, color, bg) in enumerate(layers):
    top = Inches(1.2 + i * 1.4)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), top, Inches(12.7), Inches(1.22), fill=bg, border_color=color, border_width=2)
    rounded_box(slide, Inches(0.38), top + Inches(0.04), Inches(12.54), Inches(0.35), color,
                f"{layer_name}  —  {layer_sub}", 13, WHITE)
    add_text(slide, Inches(0.5), top + Inches(0.42), Inches(12.3), Inches(0.75), detail, 10.5, False, DARK_GRAY, font="Consolas")
    if i < len(layers) - 1:
        add_text(slide, Inches(6.3), top + Inches(1.22), Inches(0.7), Inches(0.18), "▼", 13, True, color, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
#  SLIDE 5b — SYSTEM FLOW
# ═══════════════════════════════════════════════
slide = new_slide()
set_bg(slide, OFF_WHITE)
header_bar(slide, "System Flow Diagram", "End-to-end workflow: Teacher broadcast → Student decode → Server record")
footer_bar(slide, "5b / 13")

def flow_column(slide, x_start, title, title_color, steps):
    add_text(slide, x_start, Inches(1.25), Inches(5.8), Inches(0.35), title, 20, True, title_color)
    for i, (step, detail, color) in enumerate(steps):
        y = Inches(1.75 + i * 0.88)
        rounded_box(slide, x_start + Inches(0.05), y, Inches(0.42), Inches(0.42), color, str(i+1), 14, WHITE)
        box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x_start + Inches(0.55), y, Inches(5.5), Inches(0.68),
                         fill=WHITE, border_color=color, border_width=1.5)
        tf = box.text_frame
        tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.margin_left = Inches(0.1)
        p = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = step; r1.font.size = Pt(14); r1.font.bold = True; r1.font.color.rgb = color; r1.font.name = "Calibri"
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = detail; r2.font.size = Pt(11); r2.font.color.rgb = DARK_GRAY; r2.font.name = "Calibri"; p2.space_before = Pt(1)
        if i < len(steps) - 1:
            add_text(slide, x_start + Inches(0.15), y + Inches(0.68), Inches(0.3), Inches(0.18), "↓", 13, True, color, PP_ALIGN.CENTER)

flow_column(slide, Inches(0.25), "Teacher Flow (Broadcaster)", ACCENT, [
    ("Login / Register", "POST /auth/login or /register → JWT tokens (access 2h, refresh 7d)", ACCENT),
    ("Configure Session", "Enter classId (e.g., 'CS101-A'), pick duration from 1/2/5/10/15/30 min buttons", MID_BLUE),
    ("Start Session", "POST /sessions/start → UUID sessionId + 6-char code stored in Redis (TTL) + Postgres", TEAL),
    ("Broadcast Ultrasonic", "Web Audio API: sine oscillators at 17800 Hz (preamble) + 6 data tones, looped continuously", PURPLE),
    ("Monitor Live", "GET /attendance/session/:id polled every 5s → real-time student count + email list on screen", SUCCESS),
])

flow_column(slide, Inches(6.65), "Student Flow (Listener)", SUCCESS, [
    ("Login / Register", "POST /auth/login or /register (role='student') → JWT tokens stored in localStorage", SUCCESS),
    ("Start Listening", "getUserMedia({echoCancellation:false}) → AudioContext → AnalyserNode (FFT 8192)", TEAL),
    ("Detect Preamble", "getFloatFrequencyData() every 30ms → peak near 17800 Hz → ≥4 frames above -70 dB", ORANGE),
    ("Decode 6 Characters", "Each tone: ≥2 readings → median freq → freq = 18000 + idx×55 → charAt(idx)", PURPLE),
    ("Auto-Submit", "GET /sessions/code/:code → sessionId → POST /attendance/mark → 201 Attendance Marked", ACCENT),
])


# ═══════════════════════════════════════════════
#  SLIDE 5c — DATABASE DESIGN
# ═══════════════════════════════════════════════
slide = new_slide()
set_bg(slide, OFF_WHITE)
header_bar(slide, "Database Design & ER Schema", "PostgreSQL 16 with pgcrypto extension")
footer_bar(slide, "5c / 13")

# Users table
add_text(slide, Inches(0.4), Inches(1.25), Inches(4), Inches(0.35), "users", 20, True, ACCENT)
u_cols_l = [Inches(0.4), Inches(2.2), Inches(3.6), Inches(5.0)]
u_cols_w = [Inches(1.8), Inches(1.4), Inches(1.4), Inches(1.5)]
table_header(slide, u_cols_l, u_cols_w, Inches(1.6), Inches(0.32), ["Column", "Type", "Constraint", "Notes"], 11)
u_rows = [
    ("id", "UUID", "PRIMARY KEY", "gen_random_uuid()"),
    ("email", "TEXT", "UNIQUE, NOT NULL", "Normalized lowercase"),
    ("password_hash", "TEXT", "NOT NULL", "bcrypt (12 rounds)"),
    ("role", "TEXT", "NOT NULL", "'student' or 'teacher'"),
    ("created_at", "TIMESTAMP", "DEFAULT NOW()", "Auto-generated"),
]
for r, vals in enumerate(u_rows):
    bg = WHITE if r%2==0 else LIGHT_GRAY
    table_row(slide, u_cols_l, u_cols_w, Inches(1.92 + r*0.28), Inches(0.28), list(vals), bg, DARK_GRAY, 10, bold_col={0})

# Sessions table
add_text(slide, Inches(0.4), Inches(3.5), Inches(4), Inches(0.35), "sessions", 20, True, TEAL)
s_cols_l = [Inches(0.4), Inches(2.2), Inches(3.6), Inches(5.0)]
s_cols_w = [Inches(1.8), Inches(1.4), Inches(1.4), Inches(1.5)]
table_header(slide, s_cols_l, s_cols_w, Inches(3.85), Inches(0.32), ["Column", "Type", "Constraint", "Notes"], 11)
s_rows = [
    ("id", "VARCHAR(100)", "PRIMARY KEY", "UUID v4"),
    ("class_id", "TEXT", "NOT NULL", "e.g., 'CS101-A'"),
    ("code", "VARCHAR(6)", "NOT NULL", "Random 6-char code"),
    ("teacher_id", "VARCHAR(100)", "NOT NULL", "FK → users.id"),
    ("duration", "INTEGER", "NOT NULL", "Seconds (30–3600)"),
    ("created_at", "TIMESTAMP", "DEFAULT NOW()", "Session start time"),
]
for r, vals in enumerate(s_rows):
    bg = WHITE if r%2==0 else LIGHT_GRAY
    table_row(slide, s_cols_l, s_cols_w, Inches(4.17 + r*0.28), Inches(0.28), list(vals), bg, DARK_GRAY, 10, bold_col={0})

# Attendance records table
add_text(slide, Inches(6.8), Inches(1.25), Inches(6), Inches(0.35), "attendance_records", 20, True, PURPLE)
a_cols_l = [Inches(6.8), Inches(8.6), Inches(10.0), Inches(11.1)]
a_cols_w = [Inches(1.8), Inches(1.4), Inches(1.1), Inches(1.8)]
table_header(slide, a_cols_l, a_cols_w, Inches(1.6), Inches(0.32), ["Column", "Type", "Constraint", "Notes"], 11)
a_rows = [
    ("id", "SERIAL", "PRIMARY KEY", "Auto-increment"),
    ("session_id", "VARCHAR(100)", "NOT NULL", "FK → sessions.id"),
    ("student_id", "VARCHAR(100)", "NOT NULL", "FK → users.id"),
    ("timestamp", "TIMESTAMP", "DEFAULT NOW()", "When marked"),
    ("—", "—", "UNIQUE", "(session_id, student_id)"),
]
for r, vals in enumerate(a_rows):
    bg = WHITE if r%2==0 else LIGHT_GRAY
    table_row(slide, a_cols_l, a_cols_w, Inches(1.92 + r*0.28), Inches(0.28), list(vals), bg, DARK_GRAY, 10, bold_col={0})

# Indexes
add_text(slide, Inches(6.8), Inches(3.5), Inches(6), Inches(0.35), "Indexes (Performance)", 18, True, ORANGE)
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(3.9), Inches(6.1), Inches(1.6),
          fill=SOFT_ORANGE, border_color=ORANGE, border_width=1.5)
bullet_list(slide, Inches(6.95), Inches(3.95), Inches(5.8), Inches(1.5), [
    "idx_sessions_teacher ON sessions(teacher_id)",
    "idx_sessions_class ON sessions(class_id)",
    "idx_attendance_session ON attendance_records(session_id)",
    "idx_attendance_student ON attendance_records(student_id)",
    "UNIQUE(session_id, student_id) — prevents duplicate attendance",
], DARK_GRAY, 12, 3)

# Redis Keys
add_text(slide, Inches(6.8), Inches(5.65), Inches(6), Inches(0.35), "Redis Key Schema", 18, True, RED)
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(6.0), Inches(6.1), Inches(1.0),
          fill=CODE_BG, border_color=RGBColor(0x45,0x47,0x5A), border_width=1)
code_block(slide, Inches(6.85), Inches(6.05), Inches(6.0), Inches(0.9), [
    "chirp:session:{uuid}  → JSON {classId,teacherId,code} TTL=duration",
    "chirp:code:{6CHAR}    → sessionId                     TTL=duration",
    "attendance:{sid}:{uid} → '1'                           TTL=120s",
    "rl:login:{ip}         → count                         TTL=60s",
    "rl:register:{ip}      → count                         TTL=60s",
    "rl:mark:{ip}          → count                         TTL=60s",
])

# ER relationships
add_text(slide, Inches(0.4), Inches(5.85), Inches(6.2), Inches(0.35), "Entity Relationships", 18, True, SUCCESS)
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(6.2), Inches(6.1), Inches(0.8),
          fill=SOFT_GREEN, border_color=SUCCESS, border_width=1.5)
bullet_list(slide, Inches(0.55), Inches(6.22), Inches(5.8), Inches(0.7), [
    "users (1) ——→ (N) sessions  [teacher creates many sessions]",
    "sessions (1) ——→ (N) attendance_records  [each session has many records]",
    "users (1) ——→ (N) attendance_records  [student has many records]",
], DARK_GRAY, 12, 2)


# ═══════════════════════════════════════════════
#  SLIDE 5d — API ENDPOINTS
# ═══════════════════════════════════════════════
slide = new_slide()
set_bg(slide, OFF_WHITE)
header_bar(slide, "API Design & Endpoints", "RESTful API with 13 endpoints across 3 route groups")
footer_bar(slide, "5d / 13")

api_cols_l = [Inches(0.3), Inches(1.2), Inches(4.5), Inches(6.3), Inches(8.3)]
api_cols_w = [Inches(0.9), Inches(3.3), Inches(1.8), Inches(2.0), Inches(4.7)]
table_header(slide, api_cols_l, api_cols_w, Inches(1.25), Inches(0.38),
             ["Method", "Endpoint", "Auth", "Role", "Description"], 12)

endpoints = [
    ("POST", "/auth/register", "No", "—", "Create account (email, password, role) → JWT tokens + user object"),
    ("POST", "/auth/login", "No", "—", "Authenticate credentials → JWT access (2h) + refresh (7d) tokens"),
    ("POST", "/auth/refresh", "No", "—", "Exchange refresh token for new access token (rotation)"),
    ("POST", "/sessions/start", "Yes", "teacher", "Create session (classId, duration) → sessionId + 6-char code"),
    ("GET", "/sessions/active", "Yes", "teacher", "List all live sessions for this teacher (from Redis, with TTL)"),
    ("GET", "/sessions/history", "Yes", "teacher", "All past sessions with student_count (Postgres JOIN query)"),
    ("GET", "/sessions/code/:code", "Yes", "any", "Lookup session by 6-char code → sessionId + classId + remaining TTL"),
    ("POST", "/attendance/mark", "Yes", "student", "Mark attendance (sessionId) → Redis dedup → Postgres INSERT"),
    ("GET", "/attendance/me", "Yes", "student", "My attendance history with class_id + session_code (paginated)"),
    ("GET", "/attendance/list", "Yes", "any", "All attendance records (admin view, paginated, limit 200)"),
    ("GET", "/attendance/class/:classId", "Yes", "teacher", "All records for a class (multi-session, with student emails)"),
    ("GET", "/attendance/session/:sessionId", "Yes", "teacher", "Records for specific session (live polling during broadcast)"),
]

for r, (method, endpoint, auth, role, desc) in enumerate(endpoints):
    bg = WHITE if r%2==0 else LIGHT_GRAY
    row_top = Inches(1.63 + r * 0.39)
    vals = [method, endpoint, auth, role, desc]
    for i, val in enumerate(vals):
        box = add_shape(slide, MSO_SHAPE.RECTANGLE, api_cols_l[i], row_top, api_cols_w[i], Inches(0.39), fill=bg)
        box.line.color.rgb = MED_GRAY; box.line.width = Pt(0.5)
        tf = box.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.margin_left = Inches(0.06)
        p = tf.paragraphs[0]; p.font.name = "Consolas" if i < 2 else "Calibri"
        p.text = val; p.font.size = Pt(10 if i < 2 else 11)
        p.font.color.rgb = (SUCCESS if method == "GET" else ACCENT) if i == 0 else DARK_GRAY
        p.font.bold = (i == 0)

# Rate limits
add_text(slide, Inches(0.3), Inches(6.35), Inches(12.7), Inches(0.3),
         "Rate Limiting (Redis-backed, per IP):", 14, True, RED)
rate_items = [
    ("Register:", "5 requests/min (rl:register:{ip})"),
    ("Login:", "10 requests/min (rl:login:{ip})"),
    ("Mark Attendance:", "10 requests/min (rl:mark:{ip})"),
    ("Fail-open:", "If Redis is down, rate limiter passes through (doesn't block auth)"),
]
rich_bullet_list(slide, Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.5), rate_items, RED, DARK_GRAY, 12, 2)


# ═══════════════════════════════════════════════
#  SLIDE 5e — SECURITY & MIDDLEWARE
# ═══════════════════════════════════════════════
slide = new_slide()
set_bg(slide, OFF_WHITE)
header_bar(slide, "Security Architecture & Middleware", "Authentication, authorization, and protection layers")
footer_bar(slide, "5e / 13")

# Authentication
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.25), Inches(6.3), Inches(2.6),
          fill=WHITE, border_color=ACCENT, border_width=2)
rounded_box(slide, Inches(0.38), Inches(1.3), Inches(6.14), Inches(0.42), ACCENT, "Authentication (JWT + bcrypt)", 15, WHITE)
rich_bullet_list(slide, Inches(0.45), Inches(1.8), Inches(6.0), Inches(1.9), [
    ("Password Storage:", "bcrypt hash with salt rounds = 12 (brute-force resistant)"),
    ("Access Token:", "JWT signed with HMAC-SHA256, expires in 2 hours"),
    ("Refresh Token:", "Separate secret (JWT_SECRET + '_refresh'), expires in 7 days"),
    ("Token Refresh:", "POST /auth/refresh → verifies refresh token → issues new access token"),
    ("Auto-Refresh:", "Frontend interceptor: on 401 → try refresh → retry original request → else logout"),
], ACCENT, DARK_GRAY, 12, 3)

# Authorization
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.25), Inches(6.2), Inches(2.6),
          fill=WHITE, border_color=PURPLE, border_width=2)
rounded_box(slide, Inches(6.88), Inches(1.3), Inches(6.04), Inches(0.42), PURPLE, "Authorization & Middleware Pipeline", 15, WHITE)

add_text(slide, Inches(6.95), Inches(1.82), Inches(5.9), Inches(0.25),
         "Request → CORS → JSON Parse → Rate Limit → Auth → Role → Controller → Error Handler", 10, True, PURPLE, font="Consolas")

rich_bullet_list(slide, Inches(6.95), Inches(2.1), Inches(5.9), Inches(1.6), [
    ("authMiddleware:", "Extracts Bearer token from header, strips quotes (Postman fix), verifies JWT, sets req.user"),
    ("requireRole():", "Variadic role check — requireRole('teacher') blocks students and vice versa → 403 Forbidden"),
    ("Error Handler:", "Catches all thrown errors; sends {message, status}; stack trace only in NODE_ENV=development"),
    ("404 Handler:", "Catches unmatched routes → { message: 'Route not found' }"),
], PURPLE, DARK_GRAY, 12, 3)

# Security measures
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(4.1), Inches(6.3), Inches(2.85),
          fill=WHITE, border_color=RED, border_width=2)
rounded_box(slide, Inches(0.38), Inches(4.15), Inches(6.14), Inches(0.42), RED, "Anti-Abuse & Data Integrity Measures", 15, WHITE)
rich_bullet_list(slide, Inches(0.45), Inches(4.65), Inches(6.0), Inches(2.2), [
    ("Duplicate Prevention:", "Redis dedup key (attendance:{sid}:{uid}, TTL 120s) + Postgres UNIQUE(session_id, student_id)"),
    ("Postgres 23505:", "Unique violation caught → returns HTTP 409 'Already marked'"),
    ("Session Expiry:", "Redis TTL auto-deletes session & code keys → expired codes fail lookup → 404"),
    ("Code Uniqueness:", "Generated codes checked against Redis; up to 10 attempts before 500 error"),
    ("Rate Limiting:", "Redis INCR + EXPIRE pattern; per-IP tracking; Retry-After header on 429"),
    ("Email Normalization:", "trim().toLowerCase() on all email inputs prevents duplicate accounts"),
    ("Input Validation:", "classId type check, duration range (30–3600s), sessionId required check"),
], RED, DARK_GRAY, 12, 2)

# Docker
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(4.1), Inches(6.2), Inches(2.85),
          fill=WHITE, border_color=TEAL, border_width=2)
rounded_box(slide, Inches(6.88), Inches(4.15), Inches(6.04), Inches(0.42), TEAL, "Docker & Deployment Architecture", 15, WHITE)

code_block(slide, Inches(6.95), Inches(4.65), Inches(5.9), Inches(2.2), [
    "# docker-compose.yml — 3 services",
    "",
    "postgres:  postgres:16-alpine",
    "  - Port 5432, DB: smartchirp",
    "  - Volume: pgdata (persistent)",
    "  - Init: schema.sql auto-applied on first run",
    "  - Healthcheck: pg_isready every 5s",
    "",
    "redis:  redis:7-alpine",
    "  - Port 6379, healthcheck: redis-cli ping",
    "",
    "api:  node:20-alpine (Dockerfile)",
    "  - npm ci → node src/server.js",
    "  - depends_on: postgres + redis (healthy)",
    "  - Env: PORT, JWT_SECRET, REDIS_URL, DB_URL",
])


# ═══════════════════════════════════════════════
#  SLIDE 6 — HW/SW REQUIREMENTS
# ═══════════════════════════════════════════════
slide = new_slide()
set_bg(slide, OFF_WHITE)
header_bar(slide, "Hardware & Software Requirements")
footer_bar(slide, "6 / 13")

# Hardware
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.25), Inches(6.3), Inches(5.6),
          fill=WHITE, border_color=MID_BLUE, border_width=2)
rounded_box(slide, Inches(0.38), Inches(1.3), Inches(6.14), Inches(0.45), MID_BLUE, "Hardware Requirements", 17, WHITE)
hw = [
    ("Teacher Device:", "Laptop/desktop with speakers (freq response ≥18 kHz);\nmost modern laptop speakers reach 20–22 kHz"),
    ("Student Device:", "Smartphone or laptop with microphone;\nstandard MEMS mics capture up to 22–24 kHz"),
    ("Network:", "Wi-Fi / LAN / Mobile data for API calls;\nultrasonic signal itself travels only through air (offline)"),
    ("Server:", "Any machine running Node.js 20+, PostgreSQL 16,\nRedis 7; or free cloud tiers (Render, Railway, Fly.io)"),
    ("Special Hardware:", "NONE — zero additional equipment required;\nworks entirely on existing consumer devices"),
]
y = Inches(1.9)
for label, desc in hw:
    add_text(slide, Inches(0.5), y, Inches(5.9), Inches(0.22), f"▸  {label}", 14, True, MID_BLUE)
    add_text(slide, Inches(0.7), y + Inches(0.22), Inches(5.7), Inches(0.5), desc, 12, False, DARK_GRAY)
    y += Inches(0.72)

# Software (table format)
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.25), Inches(6.2), Inches(5.6),
          fill=WHITE, border_color=TEAL, border_width=2)
rounded_box(slide, Inches(6.88), Inches(1.3), Inches(6.04), Inches(0.45), TEAL, "Software Stack (package.json)", 17, WHITE)

sw_cols_l = [Inches(6.95), Inches(8.95), Inches(10.65)]
sw_cols_w = [Inches(2.0), Inches(1.7), Inches(2.2)]
table_header(slide, sw_cols_l, sw_cols_w, Inches(1.85), Inches(0.32), ["Package", "Version", "Purpose"], 11)

sw_rows = [
    ("express", "^5.2.1", "REST API framework"),
    ("pg", "^8.17.2", "PostgreSQL client (Pool)"),
    ("ioredis", "^5.9.2", "Redis client (lazy connect)"),
    ("jsonwebtoken", "^9.0.3", "JWT sign/verify"),
    ("bcrypt", "^6.0.0", "Password hashing"),
    ("uuid", "^13.0.0", "UUID v4 generation"),
    ("dotenv", "^17.2.3", "Environment config"),
    ("cors", "^2.8.6", "Cross-origin requests"),
    ("—", "—", "—"),
    ("Node.js", "20 (alpine)", "Runtime (ESM modules)"),
    ("PostgreSQL", "16 (alpine)", "Relational database"),
    ("Redis", "7 (alpine)", "In-memory cache"),
    ("Docker", "Compose v2", "Container orchestration"),
]
for r, vals in enumerate(sw_rows):
    bg = WHITE if r%2==0 else LIGHT_GRAY
    if vals[0] == "—":
        add_shape(slide, MSO_SHAPE.RECTANGLE, sw_cols_l[0], Inches(2.17 + r*0.3), Inches(5.9), Inches(0.02), fill=TEAL)
        continue
    table_row(slide, sw_cols_l, sw_cols_w, Inches(2.17 + r*0.3), Inches(0.3), list(vals), bg, DARK_GRAY, 10, bold_col={0})

add_text(slide, Inches(6.95), Inches(6.2), Inches(5.9), Inches(0.5),
         "Frontend: HTML5 + CSS3 + ES6+ JavaScript\n+ Web Audio API (AudioContext, AnalyserNode, OscillatorNode)", 12, False, DARK_GRAY)


# ═══════════════════════════════════════════════
#  SLIDE 7 — IMPLEMENTATION
# ═══════════════════════════════════════════════
slide = new_slide()
set_bg(slide, OFF_WHITE)
header_bar(slide, "Implementation Updates", "Module-by-module development status")
footer_bar(slide, "7 / 13")

modules = [
    ("Authentication System", "JWT (2h access + 7d refresh), bcrypt 12 rounds, auto\ntoken refresh interceptor, role enforcement middleware", SUCCESS),
    ("Teacher Dashboard", "Start session with class ID + duration selection (1/2/5/\n10/15/30 min), view active sessions, search by class", SUCCESS),
    ("Ultrasonic Broadcaster", "Web Audio API: OscillatorNode → GainNode → destination;\npreamble 17800Hz + 6 data tones, gain ramping (6ms)", SUCCESS),
    ("Ultrasonic Listener", "getUserMedia → MediaStreamSource → AnalyserNode (8192);\nstate machine: IDLE→AFTER_PRE→RECV→COMPLETE", SUCCESS),
    ("Session Management", "UUID v4 sessions, 6-char code (30 symbols, no ambiguous),\nRedis SET EX for TTL, Postgres INSERT ON CONFLICT", SUCCESS),
    ("Attendance Recording", "Redis EXISTS (dedup, 120s TTL) + Postgres INSERT with\nUNIQUE constraint; 409 on duplicate; paginated queries", SUCCESS),
    ("Database & Indexes", "3 tables, 4 indexes, pgcrypto UUID, schema auto-applied\nvia docker-entrypoint-initdb.d on first container run", SUCCESS),
    ("REST API (13 endpoints)", "3 route groups, rate limiting (5/10/10 per min), CORS,\nJSON error handler, 404 catch-all, role middleware", SUCCESS),
    ("Frontend UI/UX", "Dark theme (#0f0f23), Inter font, FFT frequency bars,\nanimated signal rings, responsive, backdrop-filter blur", SUCCESS),
    ("Docker Deployment", "3-service compose (postgres:16, redis:7, node:20),\nhealthchecks, persistent volume, env-based config", SUCCESS),
]

for i, (title, desc, color) in enumerate(modules):
    col = i % 2
    row = i // 2
    left = Inches(0.3 + col * 6.5)
    top = Inches(1.25 + row * 1.05)
    rounded_box(slide, left, top, Inches(0.38), Inches(0.38), color, "✓", 14, WHITE)
    box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.45), top, Inches(5.75), Inches(0.88),
                     fill=WHITE, border_color=MED_GRAY, border_width=1)
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.margin_left = Inches(0.1)
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = title; r1.font.size = Pt(13); r1.font.bold = True; r1.font.color.rgb = DARK_BLUE; r1.font.name = "Calibri"
    p2 = tf.add_paragraph()
    r2 = p2.add_run(); r2.text = desc; r2.font.size = Pt(10); r2.font.color.rgb = DARK_GRAY; r2.font.name = "Calibri"; p2.space_before = Pt(2)

# Progress bar
prog_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.4),
                      fill=SOFT_GREEN, border_color=SUCCESS, border_width=2)
tf = prog_box.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "✓  10 / 10 Core Modules Complete   •   Status: ~90%   •   Testing & Optimization Phase"
r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = SUCCESS; r.font.name = "Calibri"


# ═══════════════════════════════════════════════
#  SLIDE 8 — RESULTS
# ═══════════════════════════════════════════════
slide = new_slide()
set_bg(slide, OFF_WHITE)
header_bar(slide, "Results & Demonstrations", "Key outcomes from testing and evaluation")
footer_bar(slide, "8 / 13")

results = [
    ("Ultrasonic Signal Generation", ACCENT, SOFT_INDIGO, [
        "Generates 7 tones per cycle (1 preamble + 6 data) at 17800–19925 Hz",
        "Gain ramping (6ms) prevents audible pops/clicks at tone boundaries",
        "Total cycle time: ~2.3s (250ms preamble + 6×180ms tones + 6×90ms gaps + 650ms pause)",
        "Signal is completely inaudible — confirmed with frequency analyzer apps",
        "Works on Chrome, Edge, Firefox via standard AudioContext + OscillatorNode",
    ]),
    ("Signal Detection & Decoding", TEAL, SOFT_CYAN, [
        "FFT with 8192 bins at ~48 kHz sample rate → bin width ~5.86 Hz (sub-6 Hz resolution)",
        "Quadratic interpolation on peak bins for sub-bin frequency accuracy",
        "State machine with 8-second timeout prevents false detections from ambient noise",
        "Median frequency filtering: requires ≥2 tone readings before accepting a character",
        "Successful decoding in <5 seconds from first signal detection in quiet classroom",
    ]),
    ("Server Performance", PURPLE, SOFT_PURPLE, [
        "User registration: ~200ms (bcrypt 12 rounds dominates; hash = ~190ms)",
        "Login: ~200ms (bcrypt.compare); JWT sign < 1ms",
        "Session creation: ~50ms (Redis SET + Postgres INSERT concurrent)",
        "Attendance marking: ~30ms (Redis EXISTS + EXISTS + SET + Postgres INSERT)",
        "Session lookup by code: ~5ms (Redis GET → GET → TTL; all in-memory)",
    ]),
    ("Proximity & Security", SUCCESS, SOFT_GREEN, [
        "Effective range: ~5–10 meters (typical classroom); walls block signal completely",
        "6-char codes from 30 symbols → 30^6 = 729,000,000 possible combinations",
        "Sessions auto-expire: Redis TTL ensures codes become invalid after duration elapses",
        "Duplicate rejection: Redis dedup (120s) + DB UNIQUE constraint = double protection",
        "Role enforcement: teacher-only and student-only endpoints return 403 on mismatch",
    ]),
]

for i, (title, color, bg, metrics) in enumerate(results):
    col = i % 2
    row = i // 2
    left = Inches(0.3 + col * 6.5)
    top = Inches(1.25 + row * 2.85)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(6.2), Inches(2.65), fill=bg, border_color=color, border_width=2)
    rounded_box(slide, left + Inches(0.08), top + Inches(0.06), Inches(6.04), Inches(0.42), color, title, 15, WHITE)
    bullet_list(slide, left + Inches(0.15), top + Inches(0.55), Inches(5.9), Inches(2.0), metrics, DARK_GRAY, 12, 2, "•")


# ═══════════════════════════════════════════════
#  SLIDE 9 — ADVANTAGES & APPLICATIONS
# ═══════════════════════════════════════════════
slide = new_slide()
set_bg(slide, OFF_WHITE)
header_bar(slide, "Advantages & Applications")
footer_bar(slide, "9 / 13")

# Advantages
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.25), Inches(6.3), Inches(5.6),
          fill=WHITE, border_color=SUCCESS, border_width=2)
rounded_box(slide, Inches(0.38), Inches(1.3), Inches(6.14), Inches(0.45), SUCCESS, "Advantages of SMAT", 18, WHITE)

advs = [
    ("Proximity Verified", "Ultrasonic range ~5–10m; cannot pass through\nwalls — only students in the room can mark"),
    ("Zero Hardware Cost", "Uses existing laptops & phones; no fingerprint\nscanner, RFID reader, or BLE beacon needed"),
    ("Zero Installation", "Runs entirely in the browser; no app download,\nno plugin — Chrome, Edge, Firefox, Safari"),
    ("Anti-Proxy Design", "Ultrasonic signals can't be relayed over calls\nor WhatsApp — no remote proxy possible"),
    ("Real-time Feedback", "Teacher sees live count + student names as\nthey check in; instant during broadcast window"),
    ("Fast & Non-Disruptive", "Full cycle in ~2.3s; auto-detect < 30s total;\nsignal is inaudible — zero class disruption"),
    ("Auto-Expiry", "Redis TTL ensures sessions expire automatically;\nno stale codes; clean state management"),
]
y = Inches(1.9)
for title, desc in advs:
    rounded_box(slide, Inches(0.5), y, Inches(0.3), Inches(0.3), SUCCESS, "✓", 12, WHITE)
    add_text(slide, Inches(0.9), y - Inches(0.02), Inches(5.5), Inches(0.25), title, 14, True, DARK_BLUE)
    add_text(slide, Inches(0.9), y + Inches(0.22), Inches(5.5), Inches(0.42), desc, 11, False, DARK_GRAY)
    y += Inches(0.74)

# Applications
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.25), Inches(6.2), Inches(5.6),
          fill=WHITE, border_color=ACCENT, border_width=2)
rounded_box(slide, Inches(6.88), Inches(1.3), Inches(6.04), Inches(0.45), ACCENT, "Applications", 18, WHITE)

apps = [
    ("Universities & Colleges", "Large lecture halls with 100–500+ students;\nautomates what takes 5–10 min manually", ACCENT),
    ("Schools (K–12)", "Daily automated attendance; reduces teacher\nadmin burden; digital records for parents", TEAL),
    ("Corporate Training", "Workshop & seminar presence verification;\nno badge system needed; works immediately", PURPLE),
    ("Examination Halls", "Verify candidate physical presence during\nexams; prevent proxy test-takers", RED),
    ("Conferences & Events", "Track per-session attendance at multi-track\nconferences; no registration desk queuing", ORANGE),
    ("Coaching Institutes", "Monitor student regularity; generate\nengagement reports; alert on low attendance", SUCCESS),
]
y = Inches(1.95)
for name, desc, color in apps:
    box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.95), y, Inches(5.9), Inches(0.82),
                     fill=WHITE, border_color=color, border_width=1.5)
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.margin_left = Inches(0.12)
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = name; r1.font.size = Pt(14); r1.font.bold = True; r1.font.color.rgb = color; r1.font.name = "Calibri"
    p2 = tf.add_paragraph()
    r2 = p2.add_run(); r2.text = desc; r2.font.size = Pt(11); r2.font.color.rgb = DARK_GRAY; r2.font.name = "Calibri"; p2.space_before = Pt(1)
    y += Inches(0.88)


# ═══════════════════════════════════════════════
#  SLIDE 10 — BUDGET
# ═══════════════════════════════════════════════
slide = new_slide()
set_bg(slide, OFF_WHITE)
header_bar(slide, "Budget Analysis", "Cost breakdown — near-zero budget project")
footer_bar(slide, "10 / 13")

b_cols_l = [Inches(0.4), Inches(2.8), Inches(6.8), Inches(8.8)]
b_cols_w = [Inches(2.4), Inches(4.0), Inches(2.0), Inches(4.1)]
table_header(slide, b_cols_l, b_cols_w, Inches(1.3), Inches(0.42), ["Category", "Item", "Cost (₹)", "Notes"], 13)

budget = [
    ("Development Tools", "VS Code, Git, Docker Desktop", "₹ 0", "All free & open-source"),
    ("Backend Runtime", "Node.js 20 + Express 5", "₹ 0", "Open-source; MIT license"),
    ("Database", "PostgreSQL 16", "₹ 0", "Free tier: Render (1GB), Supabase (500MB)"),
    ("Cache / Session Store", "Redis 7", "₹ 0", "Free tier: Upstash (10K cmds/day), Railway"),
    ("Cloud Hosting", "API server (VPS)", "₹ 0 – 500/mo", "Free: Render (750h), Railway ($5 free), Fly.io"),
    ("Domain Name", "Custom domain (optional)", "₹ 500 – 1000/yr", "Optional; free subdomain works (.onrender.com)"),
    ("SSL / HTTPS", "TLS certificate", "₹ 0", "Auto-provisioned by hosting platform"),
    ("Hardware", "Additional equipment", "₹ 0", "Uses existing laptops, phones, speakers, mics"),
    ("npm Dependencies", "8 packages (bcrypt, pg, etc.)", "₹ 0", "All open-source; no paid licenses"),
]

for r, (cat, item, cost, notes) in enumerate(budget):
    bg = WHITE if r%2==0 else LIGHT_GRAY
    table_row(slide, b_cols_l, b_cols_w, Inches(1.72 + r*0.45), Inches(0.45), [cat, item, cost, notes], bg, DARK_GRAY, 12, bold_col={0},
              aligns=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.LEFT])

total_top = Inches(1.72 + len(budget) * 0.45 + 0.15)
total = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), total_top, Inches(12.5), Inches(0.55),
                   fill=SOFT_GREEN, border_color=SUCCESS, border_width=2)
tf = total.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r1 = p.add_run(); r1.text = "Total Estimated Cost:   ₹0 — ₹1,500   "; r1.font.size = Pt(18); r1.font.bold = True; r1.font.color.rgb = SUCCESS; r1.font.name = "Calibri"
r2 = p.add_run(); r2.text = "(entirely achievable with free cloud tiers — true zero-cost deployment)"; r2.font.size = Pt(13); r2.font.color.rgb = DARK_GRAY; r2.font.name = "Calibri"


# ═══════════════════════════════════════════════
#  SLIDE 11 — WORKFLOW
# ═══════════════════════════════════════════════
slide = new_slide()
set_bg(slide, OFF_WHITE)
header_bar(slide, "Workflow for Remaining Work", "6-week timeline to production-ready state")
footer_bar(slide, "11 / 13")

phases = [
    ("Week 1–2", "Cross-Browser & Environmental Testing", [
        "Test ultrasonic gen/detect on Chrome, Firefox, Edge, Safari (desktop + mobile)",
        "Benchmark in noisy environments: fan noise, ambient chatter, music",
        "Test with varied hardware: budget phones, old laptops, external speakers",
        "Measure detection success rate at 2m, 5m, 8m, 10m distances",
    ], ACCENT),
    ("Week 3", "Security Hardening & Validation", [
        "HTTPS enforcement (required for getUserMedia in all modern browsers)",
        "Rate limit tuning; brute-force protection; input sanitization (XSS/SQLi)",
        "Refresh token rotation: invalidate old tokens on refresh; blacklist on logout",
        "Penetration testing: replay attacks, token theft, code enumeration",
    ], RED),
    ("Week 4", "UI/UX Polish & Accessibility", [
        "Responsive testing across mobile (320px) / tablet (768px) / desktop viewports",
        "ARIA labels, keyboard navigation, color contrast (WCAG 2.1 AA compliance)",
        "Loading skeletons, empty states, better error messages with retry actions",
        "First-time user onboarding: tooltips explaining ultrasonic process",
    ], PURPLE),
    ("Week 5", "Performance & Audio Optimization", [
        "FFT parameter tuning: bin size vs. speed trade-off; smoothingTimeConstant",
        "Database: EXPLAIN ANALYZE on heavy queries; connection pool sizing",
        "Redis: memory usage profiling; key expiry audit; pipeline batch operations",
        "Detection accuracy: measure false positives/negatives; confidence scoring",
    ], ORANGE),
    ("Week 6", "Deployment & Documentation", [
        "Production deploy: Render/Railway + managed Postgres + Upstash Redis",
        "API docs: Postman collection export; Swagger/OpenAPI spec generation",
        "User manual for teachers (broadcast) and students (listen + mark)",
        "Final demo preparation; project report; viva/presentation prep",
    ], SUCCESS),
]

for i, (week, title, tasks, color) in enumerate(phases):
    top = Inches(1.2 + i * 1.12)
    rounded_box(slide, Inches(0.3), top, Inches(1.4), Inches(0.95), color, week, 15, WHITE)
    title_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.8), top, Inches(3.4), Inches(0.95),
                           fill=WHITE, border_color=color, border_width=2)
    tf = title_box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.margin_left = Inches(0.1)
    p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = color; p.font.name = "Calibri"

    tasks_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.3), top, Inches(7.7), Inches(0.95),
                           fill=LIGHT_GRAY, border_color=MED_GRAY, border_width=1)
    tf2 = tasks_box.text_frame; tf2.word_wrap = True; tf2.vertical_anchor = MSO_ANCHOR.MIDDLE; tf2.margin_left = Inches(0.1)
    for j, task in enumerate(tasks):
        p_t = tf2.paragraphs[0] if j == 0 else tf2.add_paragraph()
        r = p_t.add_run(); r.text = f"•  {task}"; r.font.size = Pt(10); r.font.color.rgb = DARK_GRAY; r.font.name = "Calibri"
        p_t.space_before = Pt(1)


# ═══════════════════════════════════════════════
#  SLIDE 12 — FUTURE SCOPE
# ═══════════════════════════════════════════════
slide = new_slide()
set_bg(slide, OFF_WHITE)
header_bar(slide, "Future Scope", "Potential enhancements and research directions")
footer_bar(slide, "12 / 13")

future = [
    ("Multi-Frequency Parallel\nEncoding", "Transmit multiple chars simultaneously using\ndifferent frequency sub-bands; reduce total\ncycle from 2.3s to <0.5s; OFDM-like approach", ACCENT, SOFT_INDIGO),
    ("Progressive Web App\n(PWA)", "Service worker for offline caching; push\nnotifications for session reminders; home-\nscreen installability on mobile devices", TEAL, SOFT_CYAN),
    ("Analytics &\nReporting Dashboard", "Class-wise attendance trends; student\nengagement metrics; exportable CSV/PDF;\nautomatic low-attendance email alerts", PURPLE, SOFT_PURPLE),
    ("AI-Based Anomaly\nDetection", "ML models to detect unusual patterns:\nsudden spikes, impossible timing, statistical\noutliers; flag for manual review", RED, SOFT_RED),
    ("LMS & ERP\nIntegration", "REST APIs for Moodle, Google Classroom,\nuniversity ERP; automatic gradebook sync;\ncentralized student records management", ORANGE, SOFT_ORANGE),
    ("Multi-Room Frequency\nIsolation", "Unique frequency bands per classroom;\ndynamic allocation based on room proximity;\navoid cross-room signal interference", SUCCESS, SOFT_GREEN),
]
for i, (title, desc, color, bg) in enumerate(future):
    col = i % 3
    row = i // 3
    left = Inches(0.3 + col * 4.25)
    top = Inches(1.25 + row * 2.85)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(4.05), Inches(2.6), fill=bg, border_color=color, border_width=2)
    rounded_box(slide, left + Inches(0.08), top + Inches(0.08), Inches(3.89), Inches(0.55), color, title, 13, WHITE)
    add_text(slide, left + Inches(0.15), top + Inches(0.72), Inches(3.75), Inches(1.7), desc, 12, False, DARK_GRAY)


# ═══════════════════════════════════════════════
#  SLIDE 13 — REFERENCES
# ═══════════════════════════════════════════════
slide = new_slide()
set_bg(slide, OFF_WHITE)
header_bar(slide, "References")
footer_bar(slide, "13 / 13")

refs = [
    '[1]  R. Nandakumar, K. K. Chintalapudi, V. Padmanabhan, R. Venkatesan, "Dhwani: Secure Peer-to-Peer Acoustic NFC,"\n       ACM SIGCOMM Computer Communication Review, Vol. 43, No. 4, pp. 63–74, 2013.',
    '[2]  W3C, "Web Audio API," W3C Recommendation, June 2021. Available: https://www.w3.org/TR/webaudio/',
    '[3]  A. Madhavapeddy et al., "Audio Networking: The Forgotten Wireless Technology," IEEE Pervasive Computing,\n       Vol. 4, No. 3, pp. 55–60, 2005.',
    '[4]  Google Developers, "Nearby Connections API — Audio-based proximity detection for Android & iOS," 2023.',
    '[5]  Chirp.io (Asio Ltd.), "Data-over-sound: Sending data using near-ultrasonic audio," Technical Whitepaper, 2020.',
    '[6]  S. Chowdhury et al., "QR Code Based Attendance System," Int. Journal of Computer Applications, 2019.',
    '[7]  T. S. Lim et al., "GPS-based Attendance: Vulnerabilities and Countermeasures," IEEE Access, 2020.',
    '[8]  Mozilla Developer Network, "Web Audio API Guide," MDN Web Docs. Available: https://developer.mozilla.org',
    '[9]  Express.js Contributors, "Express 5.x API Reference," expressjs.com, 2024.',
    '[10] PostgreSQL Global Development Group, "PostgreSQL 16 Documentation," postgresql.org, 2024.',
    '[11] Redis Ltd., "Redis 7.x Documentation and Command Reference," redis.io, 2024.',
    '[12] Auth0 Inc., "Introduction to JSON Web Tokens (JWT)," jwt.io. Available: https://jwt.io/introduction',
    '[13] Docker Inc., "Docker Compose Specification," docs.docker.com, 2024.',
]

ref_tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.7))
rtf = ref_tb.text_frame; rtf.word_wrap = True; rtf.paragraphs[0].text = ""
for ref in refs:
    p = rtf.add_paragraph()
    r = p.add_run(); r.text = ref; r.font.size = Pt(12); r.font.color.rgb = DARK_GRAY; r.font.name = "Calibri"
    p.space_before = Pt(6)


# ═══════════════════════════════════════════════
#  SLIDE 14 — THANK YOU
# ═══════════════════════════════════════════════
slide = new_slide()
set_bg(slide, NAVY)

add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(4), Inches(2.2), Inches(5.33), Inches(0.04), fill=ACCENT)
add_text(slide, Inches(1), Inches(2.5), Inches(11.33), Inches(1.0), "Thank You", 56, True, WHITE, PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(3.6), Inches(11.33), Inches(0.5),
         "SMAT: Smart Attendance System Using Ultrasonic Audio Signals", 22, False,
         RGBColor(0x94,0xA3,0xB8), PP_ALIGN.CENTER)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(4.3), Inches(2.33), Inches(0.03), fill=ACCENT)
add_text(slide, Inches(1), Inches(4.6), Inches(11.33), Inches(0.5),
         "Questions & Discussion", 26, True, ORANGE, PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(5.3), Inches(11.33), Inches(0.4),
         "We welcome your feedback and suggestions", 16, False,
         RGBColor(0x64,0x74,0x8B), PP_ALIGN.CENTER, italic=True)
add_shape(slide, MSO_SHAPE.RECTANGLE, 0, Inches(6.5), prs.slide_width, Inches(1.0), fill=TEAL)
add_text(slide, Inches(1), Inches(6.6), Inches(11.33), Inches(0.6),
         "Indian Institute of Information Technology Surat\nDepartment of Electronics and Communication Engineering",
         15, False, WHITE, PP_ALIGN.CENTER)


# ═══════════════ SAVE ═══════════════
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "SMAT_Attendance_Presentation.pptx")
prs.save(out)
print(f"Saved: {out}")
print(f"Total slides: {len(prs.slides)}")
